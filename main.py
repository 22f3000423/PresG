from pptx import Presentation
import shutil
def parse_slides_content(slides_content):
    """
    Parses LLM response into a list of (title, content) tuples.
    Expects response in markdown or numbered format.
    """
    slides = []
    for block in slides_content.split('\n\n'):
        lines = block.strip().split('\n')
        if len(lines) >= 2:
            title = lines[0].replace('#', '').strip()
            content = '\n'.join(lines[1:]).strip()
            slides.append((title, content))
    return slides
import requests
import os
def call_llm_api(text, guidance, api_key, provider="openai"):
    """
    Calls the specified LLM provider API to split text into slide content.
    This is a placeholder for OpenAI; you can adapt for Anthropic, Gemini, etc.
    """
    if provider == "openai":
        url = "https://api.openai.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        prompt = f"Split the following text into PowerPoint slides. {guidance if guidance else ''} Text: {text}"
        data = {
            "model": "gpt-3.5-turbo",
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 1024
        }
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]
    # Add other providers here
    return ""
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import FileResponse

app = FastAPI()

@app.get("/")
def read_root():
    return {"message": "Presentation Generator API is running."}

# Endpoint for uploading text, guidance, API key, and template file
@app.post("/generate")
def generate_presentation(
    text: str = Form(...),
    guidance: str = Form(None),
    api_key: str = Form(...),
    template: UploadFile = File(...)
):
    # Ensure templates directory exists
    os.makedirs("templates", exist_ok=True)
    template_path = f"templates/{template.filename}"
    with open(template_path, "wb") as f:
        f.write(template.file.read())

    # Call LLM API to split text into slide content
    try:
        slides_content = call_llm_api(text, guidance, api_key)
    except Exception as e:
        return {"error": str(e)}

    # Parse slide content
    slides = parse_slides_content(slides_content)

    # Generate new presentation from template
    output_path = f"generated/{template.filename.replace('.pptx', '_generated.pptx')}"
    os.makedirs("generated", exist_ok=True)
    shutil.copy(template_path, output_path)
    prs = Presentation(output_path)

    # Remove all slides from template (optional, if you want a clean start)
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Add slides from parsed content
    for title, content in slides:
        slide_layout = prs.slide_layouts[0]  # Use title+content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    prs.save(output_path)

    return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=os.path.basename(output_path))
